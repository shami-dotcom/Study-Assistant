import os, json
import requests
from PyPDF2 import PdfReader
from fpdf import FPDF
from gtts import gTTS
import speech_recognition as sr
from pptx import Presentation
from pptx.util import Pt
import networkx as nx
import matplotlib.pyplot as plt
import gradio as gr

# ======= Config =======
API_KEY = os.getenv("GROQ_API_KEY")
MODEL = "llama-3.3-70b-versatile"

MEMORY_FILE = "student_memory.json"
if not os.path.exists(MEMORY_FILE):
    with open(MEMORY_FILE, "w") as f:
        json.dump({"students": {}, "global": {}}, f, indent=2)

notes_memory = {}

# ======= PDF Extraction =======
def extract_text_from_pdf(file_obj):
    try:
        reader = PdfReader(file_obj)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

# ======= Groq Chat Wrapper =======
def groq_chat(system_prompt, user_prompt):
    if not API_KEY:
        return "GROQ_API_KEY not set!"
    url = "https://api.groq.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {API_KEY}"}
    data = {
        "model": MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
    }
    try:
        resp = requests.post(url, headers=headers, json=data)
        resp.raise_for_status()
        return resp.json()["choices"][0]["message"]["content"]
    except Exception as e:
        return f"LLM Error: {str(e)}"

# ======= Core Functions =======
def summarize_text(text):
    return groq_chat("Summarize concisely with bullets.", f"Summarize:\n{text}")

def generate_mcqs(text, n=5):
    return groq_chat("Create MCQs.", f"Generate {n} MCQs with answers:\n{text}")

def generate_study_plan(summary, days=7):
    return groq_chat("Create a study plan.", f"Create {days}-day plan:\n{summary}")

def answer_from_notes(question, notes):
    return groq_chat("Answer strictly from notes.", f"Notes:\n{notes}\nQuestion: {question}")

def generate_mindmap_outline(text):
    return groq_chat("Create hierarchical outline.", f"Outline:\n{text}")

# ======= Flashcards PDF =======
def create_flashcards_pdf(mcq_text, filename="flashcards.pdf"):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for line in mcq_text.splitlines():
        pdf.multi_cell(0, 7, line)
    pdf.output(filename)
    return filename

# ======= PPT export =======
def summary_to_pptx(summary_text, filename="summary_slides.pptx"):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    bullets = [l.strip() for l in summary_text.splitlines() if l.strip()]
    for i in range(0, len(bullets), 6):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Summary"
        tf = slide.placeholders[1].text_frame
        for b in bullets[i:i + 6]:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
    prs.save(filename)
    return filename

# ======= Mindmap =======
def draw_mindmap_from_outline(outline_text, filename="mindmap.png"):
    G = nx.Graph()
    root = "Central Topic"
    G.add_node(root)
    for line in outline_text.splitlines():
        line = line.strip()
        if ":" in line:
            parent, subs = line.split(":",1)
            parent = parent.strip()
            for s in [x.strip() for x in subs.split(",") if x.strip()]:
                G.add_edge(parent, s)
            G.add_edge(root, parent)
        else:
            G.add_edge(root, line)
    plt.figure(figsize=(10,8))
    pos = nx.spring_layout(G, seed=42)
    nx.draw(G, pos, with_labels=True, node_size=1500, font_size=8)
    plt.savefig(filename)
    plt.close()
    return filename

# ======= STT & TTS =======
def stt_from_audio_file(file_path):
    r = sr.Recognizer()
    try:
        with sr.AudioFile(file_path) as src:
            audio = r.record(src)
        return r.recognize_google(audio)
    except:
        return "Speech recognition failed."

def tts_save(text, filename="answer.mp3"):
    tts = gTTS(text=text, lang="en")
    tts.save(filename)
    return filename

# ======= Full Pipeline =======
def full_pipeline(pdf_file, student_id):
    notes = extract_text_from_pdf(pdf_file)
    notes_memory[student_id] = notes
    summary = summarize_text(notes)
    mcqs = generate_mcqs(summary)
    plan = generate_study_plan(summary)
    outline = generate_mindmap_outline(summary)
    mindmap_img = draw_mindmap_from_outline(outline)
    flash_pdf = create_flashcards_pdf(mcqs)
    pptx_file = summary_to_pptx(summary)
    return summary, mcqs, plan, outline, mindmap_img, flash_pdf, pptx_file

# ======= Gradio UI =======
with gr.Blocks() as demo:
    gr.Markdown("# 📘 AI Study Assistant")
    pdf_in = gr.File(label="Upload PDF")
    student_id = gr.Textbox(label="Student ID", value="default")
    btn = gr.Button("Process PDF")
    summary_out = gr.Textbox(label="Summary")
    mcq_out = gr.Textbox(label="MCQs")
    plan_out = gr.Textbox(label="Study Plan")
    outline_out = gr.Textbox(label="Mindmap Outline")
    mind_img = gr.Image(label="Mindmap")
    flash_dl = gr.File(label="Download Flashcards")
    ppt_dl = gr.File(label="Download PPT")

    btn.click(
        full_pipeline,
        inputs=[pdf_in, student_id],
        outputs=[summary_out, mcq_out, plan_out, outline_out, mind_img, flash_dl, ppt_dl]
    )

demo.launch()
